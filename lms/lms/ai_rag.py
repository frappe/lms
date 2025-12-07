import re
from typing import List, Dict, Tuple
import os

import frappe
from bs4 import BeautifulSoup
from math import sqrt
from lms.lms.ai_utils import call_embeddings_proxy, apply_source_weight, simple_guardrail_flags
from frappe.utils import now_datetime
import requests


def _text_from_html(html: str) -> str:
    soup = BeautifulSoup(html or "", "html.parser")
    return soup.get_text(separator=" ", strip=True)


def _chunk_text_by_sentences(text: str, max_chars: int = 600) -> List[str]:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return []
    # naive sentence split
    parts = re.split(r"(?<=[.!?])\s+", text)
    chunks, buf = [], ""
    for p in parts:
        if len(buf) + len(p) + 1 <= max_chars:
            buf = (buf + " " + p).strip()
        else:
            if buf:
                chunks.append(buf)
            buf = p
    if buf:
        chunks.append(buf)
    return chunks


def _make_snippet(text: str, max_chars: int | None = None) -> str:
    if max_chars is None:
        try:
            max_chars = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_snippet_chars") or 160)
        except Exception:
            max_chars = 160
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def compute_anchor(h_path: str | None, order: int | None) -> str | None:
    """Return a stable anchor id from a heading path or fallback to section order."""
    try:
        if h_path:
            import re as _re
            return "h-" + _re.sub(r"[^a-z0-9]+", "-", (h_path or "").lower()).strip("-")
        if order is not None:
            return f"sec-{order}"
        return None
    except Exception:
        return None


def chunk_lesson(lesson_name: str, max_chars: int | None = None) -> List[Dict]:
    """Return chunks from a lesson (title + content per chunk)."""
    lesson = frappe.db.get_value(
        "Course Lesson",
        lesson_name,
        ["name", "title", "content", "body", "instructor_content", "instructor_notes"],
        as_dict=True,
    )
    if not lesson:
        return []
    title = lesson.title
    chunks: List[Dict] = []

    # Prefer EditorJS content if present; else body (HTML/Markdown already stored as HTML)
    # Determine configured chunk size
    if not max_chars:
        try:
            max_chars = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_chunk_chars") or 600)
        except Exception:
            max_chars = 600

    # Student-visible content
    if lesson.content:
        try:
            import json

            data = json.loads(lesson.content)
            order = 0
            for block in data.get("blocks", []):
                t = block.get("type")
                d = block.get("data", {})
                text = None
                h_path = None
                if t in {"paragraph", "header"}:
                    text = _text_from_html(d.get("text"))
                    if t == "header":
                        h_path = text
                elif t == "list":
                    text = _text_from_html(" ".join(d.get("items") or []))
                elif t == "quote":
                    text = _text_from_html(d.get("text"))
                if text:
                    for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
                        chunks.append({"title": title, "content": ch, "h_path": h_path, "order": order, "source_type": "Lesson"})
                        order += 1
        except Exception:
            pass

    if not chunks and lesson.body:
        text = _text_from_html(lesson.body)
        order = 0
        for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
            chunks.append({"title": title, "content": ch, "h_path": None, "order": order, "source_type": "Lesson"})
            order += 1

    # Instructor notes and content (if present)
    # EditorJS-based instructor_content
    if lesson.get("instructor_content"):
        try:
            import json
            data = json.loads(lesson.get("instructor_content"))
            order = 0
            for block in data.get("blocks", []):
                t = block.get("type")
                d = block.get("data", {})
                text = None
                h_path = None
                if t in {"paragraph", "header"}:
                    text = _text_from_html(d.get("text"))
                    if t == "header":
                        h_path = text
                elif t == "list":
                    text = _text_from_html(" ".join(d.get("items") or []))
                elif t == "quote":
                    text = _text_from_html(d.get("text"))
                if text:
                    for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
                        chunks.append({"title": title, "content": ch, "h_path": h_path, "order": order, "source_type": "Instructor"})
                        order += 1
        except Exception:
            pass

    # HTML-based instructor_notes
    if lesson.get("instructor_notes"):
        text = _text_from_html(lesson.get("instructor_notes"))
        order = 0
        for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
            chunks.append({"title": title, "content": ch, "h_path": None, "order": order, "source_type": "Instructor"})
            order += 1

    return chunks


def _get_lesson_attachments(lesson_name: str) -> List[Dict]:
    """Return File records attached to this lesson."""
    try:
        files = frappe.get_all(
            "File",
            filters={
                "attached_to_doctype": "Course Lesson",
                "attached_to_name": lesson_name,
            },
            fields=[
                "name",
                "file_url",
                "file_name",
                "is_private",
                "file_size",
                "mime_type",
            ],
            limit_page_length=500,
        )
        return files or []
    except Exception:
        return []


def _resolve_file_path(file_url: str) -> str | None:
    """Map a file_url to an absolute file path on disk for this site."""
    try:
        if not file_url:
            return None
        from frappe.utils import get_site_path
        fname = os.path.basename(file_url)
        if file_url.startswith("/private/files/"):
            return get_site_path("private", "files", fname)
        if file_url.startswith("/files/"):
            return get_site_path("public", "files", fname)
        return None
    except Exception:
        return None


def _read_text_from_file_url(file_url: str) -> str | None:
    """Best-effort read of a site's file given a URL like /files/foo.txt or /private/files/bar.md.

    Only returns text for simple text-like files. Binary formats (PDF, PPTX, DOCX) are skipped for now.
    """
    try:
        if not file_url:
            return None
        fpath = _resolve_file_path(file_url)
        # Heuristic: allow only known text extensions
        ext = os.path.splitext(os.path.basename(fpath or ""))[1].lower()
        text_exts = {".txt", ".md", ".markdown", ".csv", ".json", ".xml", ".yml", ".yaml", ".py", ".js", ".ts", ".java", ".go", ".rb"}
        if ext not in text_exts:
            return None
        if not os.path.exists(fpath):
            return None
        with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
            data = f.read()
            return data
    except Exception:
        return None


def _extract_pdf_pages(file_url: str) -> List[Dict]:
    """Extract text per page from a PDF. Returns list of {page, text}."""
    try:
        fpath = _resolve_file_path(file_url)
        if not fpath or not os.path.exists(fpath):
            return []
        try:
            import PyPDF2  # type: ignore
        except Exception:
            return []
        out = []
        with open(fpath, "rb") as fh:
            reader = PyPDF2.PdfReader(fh)
            for i, pg in enumerate(reader.pages, start=1):
                try:
                    txt = pg.extract_text() or ""
                except Exception:
                    txt = ""
                if txt.strip():
                    out.append({"page": i, "text": txt})
        return out
    except Exception:
        return []


def _extract_pptx_slides(file_url: str) -> List[Dict]:
    """Extract text per slide from a PPTX. Returns list of {slide, text}."""
    try:
        fpath = _resolve_file_path(file_url)
        if not fpath or not os.path.exists(fpath):
            return []
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return []
        prs = Presentation(fpath)
        out = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            try:
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        texts.append(shp.text)
            except Exception:
                pass
            txt = "\n".join(texts).strip()
            if txt:
                out.append({"slide": i, "text": txt})
        return out
    except Exception:
        return []


def _extract_docx_blocks(file_url: str) -> List[Dict]:
    """Extract paragraphs from DOCX. Returns list of {text}."""
    try:
        fpath = _resolve_file_path(file_url)
        if not fpath or not os.path.exists(fpath):
            return []
        try:
            import docx  # type: ignore
        except Exception:
            return []
        doc = docx.Document(fpath)
        out = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                out.append({"text": t})
        return out
    except Exception:
        return []


def chunk_lesson_attachments(lesson_name: str, max_chars: int | None = None) -> List[Dict]:
    """Create chunks from attachments linked to the lesson (PDF/PPTX/DOCX/text-like).

    Respects Attachment MIME Allowlist from settings; gracefully skips unsupported types.
    """
    if not max_chars:
        try:
            max_chars = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_chunk_chars") or 600)
        except Exception:
            max_chars = 600
    chunks: List[Dict] = []
    files = _get_lesson_attachments(lesson_name)
    # Build allowlist
    allowlist_raw = frappe.db.get_single_value("LMS Settings", "assistant_rag_attachment_mime_allowlist") or ""
    allow = [a.strip().lower() for a in (allowlist_raw.split(",") if allowlist_raw else []) if a.strip()]
    def _allowed(mtype: str, ext: str) -> bool:
        if allow:
            for pat in allow:
                if mtype.lower().startswith(pat) or (pat.startswith(".") and ext == pat):
                    return True
            return False
        # default allow common types if no allowlist set
        defaults = [
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "text/",
            "application/json",
            "text/csv",
        ]
        return any(mtype.lower().startswith(p) for p in defaults)

    stats = {"total": len(files), "processed": 0, "skipped": 0, "errors": 0, "skipped_reasons": {}}
    order = 0
    for f in files:
        url = f.get("file_url")
        title = f.get("file_name") or os.path.basename(url or "")
        added = False
        # Prefer specialized extractors first
        ext = os.path.splitext(os.path.basename(url or ""))[1].lower()
        mtype = (f.get("mime_type") or "").lower()
        if not _allowed(mtype, ext):
            stats["skipped"] += 1
            stats["skipped_reasons"][mtype or ext or "unknown"] = stats["skipped_reasons"].get(mtype or ext or "unknown", 0) + 1
            continue
        if ext == ".pdf":
            pages = _extract_pdf_pages(url)
            if not pages:
                stats["errors"] += 1
            for page in pages:
                for ch in _chunk_text_by_sentences(page.get("text") or "", max_chars=max_chars):
                    chunks.append({
                        "title": f"{title} (p{page.get('page')})",
                        "content": ch,
                        "h_path": None,
                        "order": order,
                        "source_type": "File",
                        "url": url,
                        "file_name": title,
                        "page": page.get("page"),
                    })
                    order += 1
                    added = True
        elif ext == ".pptx":
            slides = _extract_pptx_slides(url)
            if not slides:
                stats["errors"] += 1
            for sl in slides:
                for ch in _chunk_text_by_sentences(sl.get("text") or "", max_chars=max_chars):
                    chunks.append({
                        "title": f"{title} (s{sl.get('slide')})",
                        "content": ch,
                        "h_path": None,
                        "order": order,
                        "source_type": "File",
                        "url": url,
                        "file_name": title,
                        "slide": sl.get("slide"),
                    })
                    order += 1
                    added = True
        elif ext == ".docx":
            blocks = _extract_docx_blocks(url)
            if not blocks:
                stats["errors"] += 1
            for blk in blocks:
                for ch in _chunk_text_by_sentences(blk.get("text") or "", max_chars=max_chars):
                    chunks.append({
                        "title": title,
                        "content": ch,
                        "h_path": None,
                        "order": order,
                        "source_type": "File",
                        "url": url,
                        "file_name": title,
                    })
                    order += 1
                    added = True
        # Fallback: simple text reader for text-like files
        if not added:
            text = _read_text_from_file_url(url or "")
            if text:
                for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
                    chunks.append(
                        {
                            "title": title,
                            "content": ch,
                            "h_path": None,
                            "order": order,
                            "source_type": "File",
                            "url": url,
                            "file_name": title,
                        }
                    )
                    order += 1
                    added = True
        stats["processed"] += 1 if added else 0
    return chunks


def simple_retrieve(course: str, lesson: str, query: str, top_k: int | None = None) -> List[Dict]:
    """Keyword overlap retrieval over AI Knowledge Chunk for the lesson.

    Returns [{title, content, h_path, order}]
    """
    query = (query or "").lower()
    if not query:
        return []
    q_terms = set(re.findall(r"\w+", query))
    # Determine top_k default
    if not top_k:
        try:
            top_k = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_top_k") or 3)
        except Exception:
            top_k = 3

    chunks = frappe.get_all(
        "AI Knowledge Chunk",
        filters={"course": course, "lesson": lesson},
        fields=["title", "content", "h_path", "order", "source_type"],
        limit_page_length=500,
    )
    # Load weights
    try:
        settings = frappe.get_single("LMS Settings")
        weights = {
            "Lesson": float(settings.get("assistant_rag_weight_lesson") or 1.0),
            "Instructor": float(settings.get("assistant_rag_weight_instructor") or 0.95),
            "File": float(settings.get("assistant_rag_weight_file") or 0.85),
        }
    except Exception:
        weights = {"Lesson": 1.0, "Instructor": 0.95, "File": 0.85}
    scored = []
    for ch in chunks:
        terms = set(re.findall(r"\w+", (ch.get("content") or "").lower()))
        overlap = len(q_terms & terms)
        if overlap:
            wscore = apply_source_weight(overlap, ch.get("source_type"), weights)
            scored.append((wscore, ch))
    scored.sort(key=lambda x: (-x[0], x[1]["order"]))
    result = []
    for _, ch in scored[: top_k or 3]:
        ch = dict(ch)
        snip = _make_snippet(ch.get("content"))
        try:
            settings = frappe.get_single("LMS Settings")
            enable_filter = bool(settings.get("assistant_guardrail_filter_content"))
            enable_inj = bool(settings.get("assistant_guardrail_injection_checks"))
            flags = simple_guardrail_flags(snip)
            if (enable_filter and flags.get("sensitive")) or (enable_inj and flags.get("injection")):
                snip = "[snippet omitted]"
        except Exception:
            pass
        ch["snippet"] = snip
        ch["anchor"] = compute_anchor(ch.get("h_path"), ch.get("order"))
        result.append(ch)
    return result


def index_lesson(lesson_name: str, course: str | None = None) -> int:
    """(Re)build chunks for a lesson. Returns chunk count."""
    if not course:
        course = frappe.db.get_value("Course Lesson", lesson_name, "course")
    frappe.db.delete("AI Knowledge Chunk", {"lesson": lesson_name})
    chunks = chunk_lesson(lesson_name)
    # Optionally include attachments if enabled
    try:
        include_attachments = bool(frappe.db.get_single_value("LMS Settings", "assistant_rag_include_attachments") or 0)
    except Exception:
        include_attachments = False
    attach_chunks = []
    if include_attachments:
        try:
            attach_chunks = chunk_lesson_attachments(lesson_name)
        except Exception:
            frappe.log_error(frappe.get_traceback(), "rag_chunk_attachments_error")
    all_chunks = chunks + attach_chunks
    for idx, ch in enumerate(all_chunks):
        frappe.get_doc(
            {
                "doctype": "AI Knowledge Chunk",
                "course": course,
                "lesson": lesson_name,
                "source_type": ch.get("source_type") or "Lesson",
                "title": ch.get("title"),
                "content": ch.get("content"),
                "h_path": ch.get("h_path"),
                "order": ch.get("order", idx),
                "chunk_id": f"{lesson_name}-{idx}",
                "url": ch.get("url"),
            }
        ).insert(ignore_permissions=True)
    count = len(all_chunks)
    try:
        # Summarize attachment indexing outcome
        notes = "Indexed lesson content and instructor notes"
        if include_attachments:
            try:
                files = _get_lesson_attachments(lesson_name)
                total_files = len(files)
                urls_with_chunks = set([c.get("url") for c in attach_chunks if c.get("url")])
                files_with_chunks = len(urls_with_chunks)
                skipped_or_empty = max(0, total_files - files_with_chunks)
                notes = f"{notes}; attachments=on (files={total_files}, with_chunks={files_with_chunks}, skipped_or_empty={skipped_or_empty})"
            except Exception:
                notes = f"{notes}; attachments=on"
        frappe.get_doc(
            {
                "doctype": "AI Knowledge Index Run",
                "course": course,
                "lesson": lesson_name,
                "status": "Success",
                "chunk_count": count,
                "last_indexed_at": now_datetime(),
                "notes": notes,
            }
        ).insert(ignore_permissions=True)
    except Exception:
        frappe.log_error(frappe.get_traceback(), "rag_index_run_log_error")
    return count


def index_course(course: str) -> int:
    """Index all lessons for a course; returns total chunks."""
    total = 0
    # Find all lessons under the course
    chapters = frappe.get_all("Chapter Reference", filters={"parent": course}, pluck="chapter")
    for ch in chapters:
        lessons = frappe.get_all("Lesson Reference", filters={"parent": ch}, pluck="lesson")
        for lesson in lessons:
            try:
                total += index_lesson(lesson, course=course)
            except Exception:
                frappe.log_error(frappe.get_traceback(), "rag_index_course_error")
    try:
        frappe.get_doc(
            {
                "doctype": "AI Knowledge Index Run",
                "course": course,
                "status": "Success",
                "chunk_count": total,
                "last_indexed_at": now_datetime(),
                "notes": "Indexed all lessons for course",
            }
        ).insert(ignore_permissions=True)
    except Exception:
        frappe.log_error(frappe.get_traceback(), "rag_index_run_course_log_error")
    return total


def _cosine_similarity(a: List[float], b: List[float]) -> float:
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x*y for x, y in zip(a, b))
    na = sqrt(sum(x*x for x in a))
    nb = sqrt(sum(y*y for y in b))
    if na == 0 or nb == 0:
        return 0.0
    return dot / (na * nb)


def ensure_embeddings_for_lesson(lesson: str, base_url: str, api_key: str, model: str, custom_headers_json: str = None) -> int:
    """Compute embeddings for chunks lacking embeddings or with different model. Returns updated count.

    Batches requests and retries with exponential backoff on failure.
    """
    # DEBUG: Add logging to understand what's happening
    print(f"🔍 WORKER JOB STARTED: ensure_embeddings_for_lesson called with lesson='{lesson}', base_url='{base_url}', model='{model}'")
    frappe.logger().info(f"🔍 WORKER JOB STARTED: ensure_embeddings_for_lesson called with: lesson='{lesson}', base_url='{base_url}', model='{model}', api_key='{api_key[:20] if api_key else 'None'}...'")
    
    chunks = frappe.get_all(
        "AI Knowledge Chunk",
        filters={"lesson": lesson},
        fields=["name", "content", "embedding_model", "embedding"],
        limit_page_length=5000,
    )
    frappe.logger().info(f"📊 Found {len(chunks)} chunks")
    
    # A chunk needs embedding if: 1) no model set, 2) different model, or 3) model set but no actual embedding data
    to_embed = [c for c in chunks if not c.get("embedding_model") or c.get("embedding_model") != model or not c.get("embedding")]
    frappe.logger().info(f"📊 Chunks to embed: {len(to_embed)}")
    
    if chunks:
        frappe.logger().info(f"📝 Sample chunk embedding_model values: {[c.get('embedding_model') for c in chunks[:3]]}")
    
    if not to_embed:
        frappe.logger().info(f"❌ No chunks to embed - returning 0")
        return 0
    
    frappe.logger().info(f"✅ Proceeding with embedding {len(to_embed)} chunks")

    # Settings for batching and retries
    try:
        settings = frappe.get_single("LMS Settings")
        batch_size = int(settings.get("assistant_embedding_batch_size") or 64)
        max_retries = int(settings.get("assistant_embedding_max_retries") or 3)
        batch_delay_ms = int(settings.get("assistant_embedding_batch_delay_ms") or 100)
    except Exception:
        batch_size = 64
        max_retries = 3
        batch_delay_ms = 100

    updated = 0
    import time as _time

    for i in range(0, len(to_embed), batch_size):
        batch = to_embed[i : i + batch_size]
        texts = [c.get("content") or "" for c in batch]
        attempt = 0
        while attempt <= max_retries:
            res = call_embeddings_proxy(base_url, api_key, model, texts, custom_headers_json=custom_headers_json)
            print(f"🔍 Embedding API response: {res}")
            frappe.logger().info(f"🔍 Embedding API response: status_code={res.get('status_code')}, error={res.get('error')}, vectors_count={len(res.get('vectors') or [])}")
            vectors = res.get("vectors") or []
            if vectors and len(vectors) == len(batch):
                break
            attempt += 1
            backoff = (2 ** (attempt - 1)) * 0.5  # 0.5s, 1s, 2s, ...
            _time.sleep(backoff)
        if not vectors:
            # give up on this batch
            continue
        for doc, vec in zip(batch, vectors):
            try:
                frappe.logger().info(f"💾 Updating chunk {doc['name']} with {len(vec)}-dim embedding")
                
                # Validate the vector data
                if not vec or not isinstance(vec, list) or len(vec) == 0:
                    frappe.logger().error(f"❌ Invalid vector data for chunk {doc['name']}: {type(vec)} with length {len(vec) if vec else 0}")
                    continue
                
                # Use Document API instead of db.set_value for more reliable updates
                chunk_doc = frappe.get_doc("AI Knowledge Chunk", doc["name"])
                frappe.logger().info(f"📄 Loaded chunk document: {chunk_doc.name}")
                
                # Set the values
                chunk_doc.embedding = vec
                chunk_doc.embedding_model = model
                frappe.logger().info(f"📝 Set embedding ({len(vec)} dims) and model ({model})")
                
                # Save with explicit error handling
                chunk_doc.save(ignore_permissions=True)
                frappe.logger().info(f"💾 Document save completed for {doc['name']}")
                
                # Verify the save worked by reloading
                test_doc = frappe.get_doc("AI Knowledge Chunk", doc["name"])
                if test_doc.embedding_model == model:
                    frappe.logger().info(f"✅ Verified: chunk {doc['name']} now has embedding_model = {test_doc.embedding_model}")
                    updated += 1
                else:
                    frappe.logger().error(f"❌ Verification failed: chunk {doc['name']} embedding_model = {test_doc.embedding_model}")
                    
            except Exception as e:
                frappe.logger().error(f"❌ Exception updating chunk {doc['name']}: {str(e)}")
                frappe.logger().error(f"❌ Exception type: {type(e).__name__}")
                import traceback
                frappe.logger().error(f"❌ Full traceback: {traceback.format_exc()}")
                frappe.log_error(frappe.get_traceback(), "rag_embed_update_error")
        # small pacing between batches
        _time.sleep(max(0, batch_delay_ms) / 1000.0)
    
    # Ensure database changes are committed
    if updated > 0:
        frappe.db.commit()
        frappe.logger().info(f"💾 Committed {updated} embedding updates to database")
    
    frappe.logger().info(f"🏁 ensure_embeddings_for_lesson completed. Updated: {updated}")
    return updated


def retrieve_with_embeddings(course: str, lesson: str, query: str, base_url: str, api_key: str, model: str, top_k: int | None = None) -> List[Dict]:
    if not query:
        return []
    # Get all chunks with embeddings for the lesson
    chunks = frappe.get_all(
        "AI Knowledge Chunk",
        filters={"course": course, "lesson": lesson, "embedding_model": ("=", model)},
        fields=["title", "content", "h_path", "order", "embedding", "source_type"],
        limit_page_length=1000,
    )
    if not chunks:
        return []
    # Embed query
    res = call_embeddings_proxy(base_url, api_key, model, [query])
    qv = (res.get("vectors") or [None])[0]
    if not qv:
        return []
    # Determine top_k
    if not top_k:
        try:
            top_k = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_top_k") or 3)
        except Exception:
            top_k = 3
    # Load weights
    try:
        settings = frappe.get_single("LMS Settings")
        weights = {
            "Lesson": float(settings.get("assistant_rag_weight_lesson") or 1.0),
            "Instructor": float(settings.get("assistant_rag_weight_instructor") or 0.95),
            "File": float(settings.get("assistant_rag_weight_file") or 0.85),
        }
    except Exception:
        weights = {"Lesson": 1.0, "Instructor": 0.95, "File": 0.85}
    scored = []
    for ch in chunks:
        vec = ch.get("embedding")
        if not vec:
            continue
        score = _cosine_similarity(qv, vec)
        wscore = apply_source_weight(score, ch.get("source_type"), weights)
        scored.append((wscore, ch))
    scored.sort(key=lambda x: (-x[0], x[1]["order"]))
    result = []
    for _, ch in scored[: top_k or 3]:
        ch = dict(ch)
        snip = _make_snippet(ch.get("content"))
        try:
            settings = frappe.get_single("LMS Settings")
            enable_filter = bool(settings.get("assistant_guardrail_filter_content"))
            enable_inj = bool(settings.get("assistant_guardrail_injection_checks"))
            flags = simple_guardrail_flags(snip)
            if (enable_filter and flags.get("sensitive")) or (enable_inj and flags.get("injection")):
                snip = "[snippet omitted]"
        except Exception:
            pass
        ch["snippet"] = snip
        ch["anchor"] = compute_anchor(ch.get("h_path"), ch.get("order"))
        result.append(ch)
    return result


def backfill_embeddings_daily():
    """Daily scheduler: enqueue embeddings for lessons missing vectors or with mismatched model.

    Uses effective config per course to resolve proxy base, key, and model. Enqueues jobs to avoid long tasks.
    """
    try:
        settings = frappe.get_single("LMS Settings")
        # Find lessons with chunks needing embeddings
        rows = frappe.db.sql(
            """
            select distinct course, lesson
            from `tabAI Knowledge Chunk`
            where coalesce(embedding_model, '') = ''
               or embedding_model != %(model)s
            """,
            {"model": settings.get("assistant_embedding_model") or ""},
            as_dict=True,
        )
        if not rows:
            return
        from lms.lms.ai_utils import determine_effective_assistant_config
        for r in rows:
            course = r.get("course")
            lesson = r.get("lesson")
            if not course or not lesson:
                continue
            course_cfg = frappe.db.get_value(
                "AI Assistant Config",
                {"course": course},
                [
                    "enable_overrides",
                    "proxy_base_url",
                    "proxy_api_key",
                    "rag_use_embeddings",
                    "embedding_model",
                ],
                as_dict=True,
            )
            eff = determine_effective_assistant_config(settings.as_dict(), course_cfg)
            if not eff.get("rag_use_embeddings"):
                continue
            base = eff.get("proxy_base_url") or ""
            key = eff.get("proxy_api_key") or ""
            model = eff.get("embedding_model") or settings.get("assistant_embedding_model")
            if not base or not model:
                continue
            try:
                frappe.enqueue(
                    ensure_embeddings_for_lesson,
                    queue="default",
                    job_name=f"rag-embed-lesson-{lesson}",
                    lesson=lesson,
                    base_url=base,
                    api_key=key,
                    model=model,
                )
            except Exception:
                frappe.log_error(frappe.get_traceback(), "rag_backfill_enqueue_error")
    except Exception:
        frappe.log_error(frappe.get_traceback(), "rag_backfill_daily_error")


def rebuild_lesson_index_and_embeddings(lesson: str, course: str | None, use_embeddings: bool, base_url: str | None, api_key: str | None, model: str | None, custom_headers_json: str = None) -> dict:
    """Rebuild chunks for a lesson then (optionally) compute embeddings.

    Returns {"chunks": int, "embedded": int}
    """
    if not course:
        course = frappe.db.get_value("Course Lesson", lesson, "course")
    chunks = index_lesson(lesson, course=course)
    embedded = 0
    if use_embeddings and base_url and model:
        try:
            embedded = ensure_embeddings_for_lesson(lesson, base_url, api_key or "", model, custom_headers_json)
        except Exception:
            frappe.log_error(frappe.get_traceback(), "rag_rebuild_embed_error")
    return {"chunks": chunks, "embedded": embedded}


def rebuild_course_index_and_embeddings(course: str, use_embeddings: bool, base_url: str | None, api_key: str | None, model: str | None, custom_headers_json: str = None) -> dict:
    """Rebuild chunks for all lessons in a course, then (optionally) compute embeddings.

    Returns {"chunks": int, "embedded": int}
    """
    total_chunks = index_course(course)
    total_embedded = 0
    if use_embeddings and base_url and model:
        import time as _time
        delay_ms = 0
        try:
            settings = frappe.get_single("LMS Settings")
            delay_ms = int(settings.get("assistant_rebuild_lesson_delay_ms") or 0)
        except Exception:
            delay_ms = 0
        chapters = frappe.get_all("Chapter Reference", filters={"parent": course}, pluck="chapter")
        for ch in chapters:
            lessons = frappe.get_all("Lesson Reference", filters={"parent": ch}, pluck="lesson")
            for lesson in lessons:
                try:
                    total_embedded += ensure_embeddings_for_lesson(lesson, base_url, api_key or "", model, custom_headers_json)
                except Exception:
                    frappe.log_error(frappe.get_traceback(), "rag_rebuild_course_embed_error")
                if delay_ms and delay_ms > 0:
                    _time.sleep(delay_ms / 1000.0)
    return {"chunks": total_chunks, "embedded": total_embedded}


def _domain_allowed(url: str) -> bool:
    try:
        from urllib.parse import urlparse
        host = urlparse(url or "").hostname or ""
        allow_raw = frappe.db.get_single_value("LMS Settings", "assistant_external_allowed_domains") or ""
        allow = [h.strip().lower() for h in allow_raw.split(",") if h.strip()]
        if not allow:
            return False
        return any(host.endswith(dom) for dom in allow)
    except Exception:
        return False


def fetch_external_to_text(url: str, timeout: int = 15) -> tuple[str | None, str | None]:
    """Fetch URL and return (mime_type, text) where text is cleaned HTML text."""
    if not _domain_allowed(url):
        return None, None
    try:
        resp = requests.get(url, timeout=timeout)
        if resp.status_code != 200:
            return None, None
        ctype = resp.headers.get("Content-Type", "").split(";")[0].strip().lower()
        text = resp.text or ""
        soup = BeautifulSoup(text, "html.parser")
        content = soup.get_text(separator=" ", strip=True)
        # Guardrails for external ingestion
        try:
            settings = frappe.get_single("LMS Settings")
            enable_filter = bool(settings.get("assistant_guardrail_filter_content"))
            enable_inj = bool(settings.get("assistant_guardrail_injection_checks"))
            flags = simple_guardrail_flags(content)
            if (enable_filter and flags.get("sensitive")) or (enable_inj and flags.get("injection")):
                try:
                    frappe.get_doc({
                        "doctype": "AI Guardrail Event",
                        "source": "external",
                        "event_type": "Sensitive" if flags.get("sensitive") else "Injection",
                        "snippet": (content or "")[:240],
                        "notes": f"Blocked external fetch for {url}",
                    }).insert(ignore_permissions=True)
                except Exception:
                    pass
                return ctype, None
        except Exception:
            pass
        return ctype, content
    except Exception:
        return None, None


def index_external_source(docname: str) -> dict:
    """Fetch and index an AI External Source document as External chunks."""
    doc = frappe.get_doc("AI External Source", docname)
    if not doc or not doc.allowed:
        return {"chunks": 0}
    mime, text = fetch_external_to_text(doc.url)
    if not text:
        frappe.db.set_value("AI External Source", docname, {"status": "Error", "error": "fetch_failed"})
        return {"chunks": 0}
    # Chunk and store
    try:
        max_chars = int(frappe.db.get_single_value("LMS Settings", "assistant_rag_chunk_chars") or 600)
    except Exception:
        max_chars = 600
    chunks = []
    order = 0
    for ch in _chunk_text_by_sentences(text, max_chars=max_chars):
        chunks.append(
            {
                "doctype": "AI Knowledge Chunk",
                "course": doc.course,
                "lesson": doc.lesson,
                "source_type": "External",
                "title": doc.title or doc.url,
                "content": ch,
                "order": order,
                "chunk_id": f"{doc.name}-{order}",
                "url": doc.url,
            }
        )
        order += 1
    for row in chunks:
        try:
            frappe.get_doc(row).insert(ignore_permissions=True)
        except Exception:
            frappe.log_error(frappe.get_traceback(), "external_chunk_insert_error")
    frappe.db.set_value(
        "AI External Source",
        docname,
        {"status": "Indexed", "last_fetched_at": now_datetime(), "mime_type": mime, "content_text": text[:5000], "error": None},
    )
    return {"chunks": len(chunks)}
